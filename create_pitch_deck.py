"""
Script to generate GSC Phase 2 PowerPoint Pitch Deck
Run: python create_pitch_deck.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_pitch_deck():
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define colors
    DARK_BG = RGBColor(15, 23, 42)  # Dark blue background
    TEAL = RGBColor(20, 184, 166)   # Teal accent
    WHITE = RGBColor(255, 255, 255)
    GRAY = RGBColor(148, 163, 184)
    
    def add_title_slide(title, subtitle=""):
        """Add a title slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = DARK_BG
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(44)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = WHITE
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Subtitle
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(9), Inches(0.8))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle
            subtitle_frame.paragraphs[0].font.size = Pt(24)
            subtitle_frame.paragraphs[0].font.color.rgb = TEAL
            subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_content_slide(title, content_items):
        """Add a content slide with bullet points"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = DARK_BG
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = TEAL
        
        # Content
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.5), Inches(5.5))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        
        for item in content_items:
            p = text_frame.add_paragraph()
            p.text = item
            p.font.size = Pt(18)
            p.font.color.rgb = WHITE
            p.level = 0
            p.space_after = Pt(10)
        
        return slide
    
    # Slide 1: Title
    add_title_slide(
        "FLOOD WATCH",
        "Community-Driven Flood Resilience System"
    )
    
    # Slide 2: Team Details
    add_content_slide("Team Details", [
        "Team Name: AIGineers",
        "Theme: Adaptation & Resilience",
        "",
        "[Add team member names, universities, and degrees here]"
    ])
    
    # Slide 3: Executive Summary
    add_content_slide("Executive Summary", [
        "Problem: 1.8B people in flood-prone areas lack early warnings",
        "Solution: WhatsApp/Telegram bot + AI verification + geospatial alerts",
        "Impact: 65% faster response (45→4.2 min), 92% user satisfaction",
        "Innovation: Accessible, AI-powered, <$0.05/user/month"
    ])
    
    # Slide 4: Phase 2 Work Summary
    add_content_slide("Phase 2 Work Summary", [
        "✅ Deployed WhatsApp & Telegram bots (multilingual)",
        "✅ Built AI model: 89.2% accuracy, 5,247 training images",
        "✅ Created admin dashboard with real-time mapping",
        "✅ Tested in 3 communities: 47 reports, 127 users, 15 volunteers",
        "✅ Results: 387 people reached, 99.2% uptime, 94% delivery rate"
    ])
    
    # Slide 5: Problem Background
    add_content_slide("Problem Background", [
        "What: No accessible early warning systems in flood areas",
        "Why: Kenya loses $300M/year, 500K+ affected annually",
        "Research: 87% had no warning, 92% own WhatsApp phones",
        "Gap: Existing systems complex or unreliable (35% failure)"
    ])
    
    # Slide 6: Proposed Solution - How It Works
    add_content_slide("How Flood Watch Works", [
        "1. REPORTING: Citizens send flood reports via WhatsApp/Telegram",
        "2. VERIFICATION: AI analyzes images (89% accuracy) + human review",
        "3. ALERTING: Geospatial system alerts users within 4.2 minutes",
        "",
        "Science: EfficientNet-B3 AI, PostGIS geofencing, multi-channel delivery"
    ])
    
    # Slide 7: Innovation & Impact
    add_content_slide("Innovation & Impact", [
        "🚀 Uses existing platforms (480M WhatsApp users in Africa)",
        "🤖 AI reduces admin workload by 70%",
        "💰 Ultra-low cost: <$0.05 per user per month",
        "🌍 Scalable: Cloud-native, serves 100K users without changes",
        "♻️ Open API for third-party integration"
    ])
    
    # Slide 8: Prototype Overview
    add_content_slide("Prototype Built", [
        "Component 1: WhatsApp & Telegram Bots (Python FastAPI)",
        "Component 2: AI Engine (EfficientNet-B3, 89.2% accuracy)",
        "Component 3: Admin Dashboard (React + PostGIS mapping)",
        "",
        "Status: Full MVP deployed, not a mockup!",
        "Uptime: 99.2% over 30 days"
    ])
    
    # Slide 9: Testing Plan & Methods
    add_content_slide("Testing Plan", [
        "Goal 1: Validate accessibility (WhatsApp vs apps)",
        "Goal 2: Test AI accuracy (target >85% in real world)",
        "Goal 3: Measure response time (target <5 minutes)",
        "Goal 4: Assess community engagement",
        "",
        "Method: 3 pilot communities (Nairobi, Homabay, Kisumu)",
        "Duration: 30 days, 127 users, 15 volunteers"
    ])
    
    # Slide 10: Testing Results
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Testing Results"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = TEAL
    
    content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.text = """Metric          Target    Result    Status
User Adoption      100       127      ✅ 127%
AI Accuracy        >85%      89.2%    ✅ 105%
Response Time      <5 min    4.2 min  ✅ 84% faster
User Satisfaction  >80%      92%      ✅ 115%
Alert Delivery     >90%      94%      ✅ 104%
System Uptime      >95%      99.2%    ✅ 104%"""
    text_frame.paragraphs[0].font.size = Pt(16)
    text_frame.paragraphs[0].font.color.rgb = WHITE
    text_frame.paragraphs[0].font.name = 'Courier New'
    
    # Slide 11: Iterations & Improvements
    add_content_slide("Iterations & Improvements", [
        "Iteration 1: Simplified bot (5→2 steps), drop-off 35%→8%",
        "Iteration 2: Added Swahili/Luo, rural engagement +55%",
        "Iteration 3: AI retraining, false positives 18%→8.3%",
        "Iteration 4: Dynamic radius (urban 2km, rural 10km)",
        "Iteration 5: Manual address search for GPS accuracy"
    ])
    
    # Slide 12: Implementation Plan - Short Term
    add_content_slide("Implementation: Next 6 Months", [
        "Goal: 15 communities across 5 counties",
        "",
        "Steps:",
        "  • Identify 12 communities (coastal, lake basin)",
        "  • Recruit 30 volunteers, conduct training",
        "  • Deploy bots, onboard 2,000+ users",
        "",
        "Cost: $45,000 (infrastructure $15K, onboarding $20K, ops $10K)"
    ])
    
    # Slide 13: Implementation Plan - Long Term
    add_content_slide("Implementation: 1-3 Years", [
        "Year 1: 20 communities, 10K users, 5 counties ($150K)",
        "Year 2: 50 communities, 50K users, 15 counties ($500K)",
        "Year 3: 100 communities, 100K users, nationwide ($1.2M)",
        "",
        "Adopters: County governments, NDMA, Kenya Red Cross, insurers",
        "",
        "Revenue: Year 3 mix - 60% gov, 30% NGO, 10% commercial"
    ])
    
    # Slide 14: Impact Assessment
    add_content_slide("Impact Assessment", [
        "Environmental: 2,500 flood events documented (Year 5)",
        "Economic: $5M losses prevented, 4:1 ROI",
        "Social: 30% mortality reduction, 50% less displacement",
        "",
        "Key Metrics:",
        "  • 100,000 people with early warnings by Year 5",
        "  • 80% reduction in emergency response time",
        "  • 100 communities with adaptive capacity"
    ])
    
    # Slide 15: Scalability Plan
    add_content_slide("Scalability Plan", [
        "Pathway 1: Geographic expansion (Kenya → East Africa)",
        "Pathway 2: Government integration (county adoption)",
        "Pathway 3: Feature expansion (drought, disease tracking)",
        "Pathway 4: Regional expansion (Uganda, Tanzania, Rwanda)",
        "",
        "Challenges: Connectivity, literacy, volunteer sustainability",
        "Solutions: SMS fallback, voice reporting, micro-stipends"
    ])
    
    # Slide 16: Thank You / Contact
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "Thank You!"
    title_frame.paragraphs[0].font.size = Pt(48)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = WHITE
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Building Resilience, One Alert at a Time 🌊🛡️\n\nTeam AIGineers\naigineers.floodwatch@gmail.com"
    subtitle_frame.paragraphs[0].font.size = Pt(20)
    subtitle_frame.paragraphs[0].font.color.rgb = TEAL
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Save presentation
    output_file = "GSC_Phase2_FloodWatch_Presentation.pptx"
    prs.save(output_file)
    print(f"PowerPoint created: {output_file}")
    print(f"Total slides: {len(prs.slides)}")
    return output_file

if __name__ == "__main__":
    try:
        create_pitch_deck()
    except ImportError:
        print("Error: python-pptx library not installed")
        print("Install it with: pip install python-pptx")
        print("\nAlternatively, you can:")
        print("1. Copy GSC_Phase2_Pitch_Deck.md to Google Docs")
        print("2. Use File -> Download -> Microsoft PowerPoint")
